#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
convert_kb.py
-------------
Build a GPT-friendly text knowledge base from Teams/SharePoint files.

Features
- Recursively read a local folder of docs (or optionally download from SharePoint via MS Graph).
- Extract text from:
    * PDF (native text + OCR fallback for diagrams/scans)
    * PPT/PPTX (slide text; basic diagram edges; Mermaid generation)
    * DOC/DOCX
    * XLS/XLSX (tabular to TSV-like text)
    * PNG/JPG (OCR)
- Chunk long outputs to stay under model token limits (configurable).
- Outputs mirror the source folder structure under /output.

NOTE: Only the "local folder" flow is required. The Graph downloader is optional.
"""

from __future__ import annotations
import os
import re
import sys
import json
import math
import shutil
from pathlib import Path
from typing import List, Tuple

# --- CONFIG (EDIT THESE) ------------------------------------------------------

CONFIG = {
    # Required: where you unzipped the Teams files
    "INPUT_DIR": r"C:\kb\raw",
    # Where to put converted text
    "OUTPUT_DIR": r"C:\kb\output",

    # Path to tesseract.exe (your user install)
    "TESSERACT_CMD": r"C:\Users\VishakhaMaheshwari\AppData\Local\Programs\Tesseract-OCR\tesseract.exe",
    # OCR language (must be installed in Tesseract)
    "OCR_LANG": "eng",

    # Token / chunk controls
    # If tiktoken is present, we tokenize; else we approximate 4 chars/token
    "MAX_TOKENS_PER_CHUNK": 7500,   # adjust to your GPT limit
    "CHUNK_OVERLAP_TOKENS": 200,

    # Optional: use MS Graph to download directly from SharePoint/Teams
    "GRAPH_DOWNLOAD": {
        "ENABLED": False,  # set True if you want to download instead of manual copy
        # You need an Entra ID app (Public client), then grant Files.Read.All, Sites.Read.All
        "TENANT_ID": "<your-tenant-id>",
        "CLIENT_ID": "<your-client-id>",  # public client (no secret)
        # Where to pull from (team site)
        "SITE_HOSTNAME": "contoso.sharepoint.com",
        "SITE_PATH": "/sites/YourTeamName",     # e.g. /sites/Engineering
        # Drive name usually "Documents" for the default doc library
        "DRIVE_NAME": "Documents",
        # Folder path inside the drive to mirror, or "" for the whole drive
        "FOLDER_PATH": "General",  # e.g. "General/Designs"
        "DOWNLOAD_DIR": r"C:\kb\raw"  # where files land locally
    },
}

# --- imports that can fail gracefully -----------------------------------------
try:
    import fitz  # PyMuPDF
except Exception as e:
    print("PyMuPDF not available; PDF text will be limited.", e)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as e:
    print("python-pptx not available; PPTX parsing disabled.", e)

try:
    import docx
except Exception as e:
    print("python-docx not available; DOCX parsing disabled.", e)

try:
    import pandas as pd
except Exception as e:
    print("pandas not available; Excel parsing disabled.", e)

try:
    from PIL import Image, ImageDraw
except Exception as e:
    print("Pillow not available; image handling limited.", e)

try:
    import pytesseract
except Exception as e:
    print("pytesseract not available; OCR disabled.", e)
    pytesseract = None

try:
    import cv2
except Exception as e:
    print("opencv not available; advanced line detection disabled.", e)
    cv2 = None

try:
    import tiktoken
except Exception:
    tiktoken = None

try:
    from tqdm import tqdm
except Exception:
    def tqdm(x, **k): return x  # fallback

# Optional Graph
try:
    import requests
    import msal
except Exception:
    requests = None
    msal = None

# --- helpers ------------------------------------------------------------------

def ensure_dirs():
    Path(CONFIG["OUTPUT_DIR"]).mkdir(parents=True, exist_ok=True)
    if CONFIG["GRAPH_DOWNLOAD"]["ENABLED"]:
        Path(CONFIG["GRAPH_DOWNLOAD"]["DOWNLOAD_DIR"]).mkdir(parents=True, exist_ok=True)

def set_tesseract_path():
    if pytesseract and CONFIG["TESSERACT_CMD"] and Path(CONFIG["TESSERACT_CMD"]).exists():
        pytesseract.pytesseract.tesseract_cmd = CONFIG["TESSERACT_CMD"]

def safe_write_text(path: Path, text: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(text)

def tokenize_len(text: str) -> int:
    if tiktoken:
        enc = tiktoken.get_encoding("cl100k_base")
        return len(enc.encode(text))
    # crude fallback: ~4 chars/token
    return max(1, math.ceil(len(text) / 4))

def chunk_text(text: str, max_tokens: int, overlap_tokens: int) -> List[str]:
    def _tok_len(s: str) -> int:
        return tokenize_len(s)

    if _tok_len(text) <= max_tokens:
        return [text]
    # tokenize precisely if available; else char-approx
    if tiktoken:
        enc = tiktoken.get_encoding("cl100k_base")
        tokens = enc.encode(text)
        chunks = []
        i = 0
        step = max(1, max_tokens - overlap_tokens)
        while i < len(tokens):
            window = tokens[i:i+max_tokens]
            chunks.append(enc.decode(window))
            i += step
        return chunks
    else:
        avg = 4
        max_chars = max_tokens * avg
        overlap_chars = overlap_tokens * avg
        chunks = []
        i = 0
        while i < len(text):
            window = text[i:i+max_chars]
            chunks.append(window)
            i += max_chars - overlap_chars
        return chunks

def clean_whitespace(s: str) -> str:
    s = re.sub(r"\r\n?", "\n", s)
    s = re.sub(r"[ \t]+\n", "\n", s)
    return s.strip()

# --- converters ---------------------------------------------------------------

def convert_pdf(path: Path) -> str:
    out = [f"[SOURCE] {path}", ""]
    if 'fitz' not in sys.modules:
        out.append("(PyMuPDF not installed)")
        return "\n".join(out)
    doc = fitz.open(path)
    for i, page in enumerate(doc, start=1):
        out.append(f"\n=== PDF Page {i} ===")
        text = page.get_text("text") or ""
        text = clean_whitespace(text)
        if text:
            out.append(text)
        # OCR fallback if page text is sparse (diagrams/scans)
        need_ocr = len(text) < 30
        if need_ocr and pytesseract:
            pix = page.get_pixmap(dpi=300, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_text = pytesseract.image_to_string(img, lang=CONFIG["OCR_LANG"])
            ocr_text = clean_whitespace(ocr_text)
            if ocr_text:
                out.append("\n[OCR]")
                out.append(ocr_text)
    doc.close()
    return "\n".join(out)

def convert_image(path: Path) -> str:
    if not pytesseract:
        return f"[SOURCE] {path}\n(OCR not available)"
    img = Image.open(path)
    txt = pytesseract.image_to_string(img, lang=CONFIG["OCR_LANG"])
    return f"[SOURCE] {path}\n{clean_whitespace(txt)}"

def convert_docx(path: Path) -> str:
    if 'docx' not in sys.modules:
        return f"[SOURCE] {path}\n(python-docx not installed)"
    d = docx.Document(str(path))
    parts = [f"[SOURCE] {path}", ""]
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    # simple tables
    for tbl in d.tables:
        parts.append("\n[Table]")
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)

def convert_excel(path: Path) -> str:
    if 'pandas' not in sys.modules:
        return f"[SOURCE] {path}\n(pandas not installed)"
    parts = [f"[SOURCE] {path}", ""]
    try:
        xls = pd.ExcelFile(path, engine=None)  # auto-detect
    except Exception as e:
        return f"[SOURCE] {path}\n(Excel read error: {e})"
    for sheet in xls.sheet_names:
        try:
            df = xls.parse(sheet)
            parts.append(f"\n=== Sheet: {sheet} ===")
            parts.append(df.to_csv(sep="\t", index=False))
        except Exception as e:
            parts.append(f"\n=== Sheet: {sheet} ===\n(Read error: {e})")
    return "\n".join(parts)

def _pptx_collect_shapes_with_text(slide):
    items = []
    for shp in slide.shapes:
        if hasattr(shp, "text_frame") and shp.text_frame:
            txt = "\n".join([p.text for p in shp.text_frame.paragraphs]).strip()
            if txt:
                items.append({
                    "id": id(shp),
                    "text": txt,
                    "left": int(shp.left), "top": int(shp.top),
                    "width": int(shp.width), "height": int(shp.height),
                    "shape_type": getattr(shp, "shape_type", None)
                })
    return items

def _pptx_guess_edges(slide, nodes):
    """Heuristic: treat connector shapes as edges and snap endpoints
    to nearest text-bearing shapes."""
    edges = []
    if not hasattr(slide, "shapes"):
        return edges

    def center(b):
        return (b["left"] + b["width"]//2, b["top"] + b["height"]//2)

    def nearest_node(x, y):
        best = None
        best_d2 = 1e18
        for n in nodes:
            cx, cy = center(n)
            d2 = (cx - x)*(cx - x) + (cy - y)*(cy - y)
            if d2 < best_d2:
                best_d2, best = d2, n
        return best

    for shp in slide.shapes:
        st = getattr(shp, "shape_type", None)
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE as _T
            is_line = st == _T.LINE or st == _T.CONNECTOR
        except Exception:
            is_line = False
        if is_line:
            x1, y1 = int(shp.left), int(shp.top)
            x2, y2 = int(shp.left + shp.width), int(shp.top + shp.height)
            s = nearest_node(x1, y1)
            t = nearest_node(x2, y2)
            if s and t and s is not t:
                edges.append((s["text"], t["text"]))
    return edges

def convert_pptx(path: Path) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Return (plain_text, edges, mermaid_blocks_per_slide)."""
    if 'pptx' not in sys.modules:
        return f"[SOURCE] {path}\n(python-pptx not installed)", [], []

    prs = Presentation(str(path))
    doc_parts = [f"[SOURCE] {path}", ""]
    mermaid_blocks = []
    all_edges = []

    for i, slide in enumerate(prs.slides, start=1):
        doc_parts.append(f"\n=== Slide {i} ===")
        nodes = _pptx_collect_shapes_with_text(slide)
        for n in nodes:
            doc_parts.append(n["text"])

        edges = _pptx_guess_edges(slide, nodes)
        all_edges.extend(edges)

        if nodes:
            labels = [n["text"].strip()[:60].replace("\n", " ") for n in nodes]
            unique_labels = []
            for lbl in labels:
                if lbl not in unique_labels:
                    unique_labels.append(lbl)

            mer = ["flowchart TD"]
            for idx, lbl in enumerate(unique_labels):
                mer.append(f"  N{idx}({json.dumps(lbl)})")
            id_by_lbl = {lbl: f"N{idx}" for idx, lbl in enumerate(unique_labels)}
            for s, t in edges:
                s_lbl = s.strip()[:60].replace("\n", " ")
                t_lbl = t.strip()[:60].replace("\n", " ")
                if s_lbl in id_by_lbl and t_lbl in id_by_lbl:
                    mer.append(f"  {id_by_lbl[s_lbl]} --> {id_by_lbl[t_lbl]}")
            mermaid_blocks.append("\n".join(mer))

    return "\n".join(doc_parts), all_edges, mermaid_blocks

# --- MS Graph (optional downloader) -------------------------------------------

GRAPH_SCOPES = ["Files.Read.All", "Sites.Read.All", "User.Read", "offline_access", "openid", "profile"]

def graph_device_login(tenant_id, client_id):
    if not (msal and requests):
        raise RuntimeError("msal/requests not installed.")
    app = msal.PublicClientApplication(client_id=client_id, authority=f"https://login.microsoftonline.com/{tenant_id}")
    flow = app.initiate_device_flow(scopes=[f"https://graph.microsoft.com/{s}" for s in GRAPH_SCOPES])
    if "user_code" not in flow:
        raise RuntimeError(f"Device flow failed: {flow}")
    print("\n=== Microsoft login ===")
    print(flow["message"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

def graph_request(token, url):
    r = requests.get(url, headers={"Authorization": f"Bearer {token}"})
    r.raise_for_status()
    return r.json()

def graph_find_site_and_drive(token, hostname, site_path, drive_name):
    site = graph_request(token, f"https://graph.microsoft.com/v1.0/sites/{hostname}:{site_path}")
    site_id = site["id"]
    drives = graph_request(token, f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives")["value"]
    drive = next((d for d in drives if d["name"].lower() == drive_name.lower()), None)
    if not drive:
        raise RuntimeError(f"Drive '{drive_name}' not found. Available: {[d['name'] for d in drives]}")
    return site_id, drive["id"]

def graph_download_folder(token, drive_id, folder_rel, dest_dir):
    base = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{folder_rel}:/children" if folder_rel else f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
    queue = [(base, Path(dest_dir))]
    while queue:
        url, outdir = queue.pop()
        data = graph_request(token, url)
        for item in data.get("value", []):
            name = item["name"]
            if item.get("folder"):
                child_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item['id']}/children"
                queue.append((child_url, outdir / name))
            else:
                outdir.mkdir(parents=True, exist_ok=True)
                dl = item["@microsoft.graph.downloadUrl"]
                resp = requests.get(dl, stream=True)
                fp = outdir / name
                with open(fp, "wb") as f:
                    shutil.copyfileobj(resp.raw, f)
                print(f"Downloaded: {fp}")
        if "@odata.nextLink" in data:
            queue.append((data["@odata.nextLink"], outdir))

# --- main processing -----------------------------------------------------------

SUPPORTED = {".pdf", ".png", ".jpg", ".jpeg", ".pptx", ".docx", ".xlsx", ".xls"}

def process_file(src: Path, base_in: Path, base_out: Path):
    rel = src.relative_to(base_in)
    stem = src.stem
    out_dir = base_out / rel.parent
    txt_out = out_dir / f"{stem}.txt"

    ext = src.suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            text = convert_pdf(src)
        elif ext in {".png", ".jpg", ".jpeg"}:
            text = convert_image(src)
        elif ext in {".docx"}:
            text = convert_docx(src)
        elif ext in {".xlsx", ".xls"}:
            text = convert_excel(src)
        elif ext in {".pptx"}:
            text, edges, mermaids = convert_pptx(src)
            if mermaids:
                for i, mer in enumerate(mermaids, start=1):
                    safe_write_text(out_dir / f"{stem}.slide{i}.mermaid.mmd", mer)
        else:
            text = f"[SOURCE] {src}\n(Unsupported extension; convert to modern type first.)"
    except Exception as e:
        text = f"[SOURCE] {src}\n(Extraction error: {e})"

    text = clean_whitespace(text)

    chunks = chunk_text(text, CONFIG["MAX_TOKENS_PER_CHUNK"], CONFIG["CHUNK_OVERLAP_TOKENS"])
    if len(chunks) == 1:
        safe_write_text(txt_out, chunks[0])
    else:
        for i, chunk in enumerate(chunks, start=1):
            safe_write_text(out_dir / f"{stem}.part{i:03d}.txt", chunk)

def main():
    ensure_dirs()
    set_tesseract_path()

    g = CONFIG["GRAPH_DOWNLOAD"]
    if g["ENABLED"]:
        if not (msal and requests):
            raise RuntimeError("Enable Graph but msal/requests not installed.")
        token = graph_device_login(g["TENANT_ID"], g["CLIENT_ID"])
        site_id, drive_id = graph_find_site_and_drive(token, g["SITE_HOSTNAME"], g["SITE_PATH"], g["DRIVE_NAME"])
        graph_download_folder(token, drive_id, g["FOLDER_PATH"], g["DOWNLOAD_DIR"])

    base_in = Path(CONFIG["INPUT_DIR"])
    base_out = Path(CONFIG["OUTPUT_DIR"])

    files = [p for p in base_in.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED]
    if not files:
        print(f"No supported files found in {base_in}.")
        return

    for p in tqdm(files, desc="Converting"):
        process_file(p, base_in, base_out)

    print("\nDone. Converted knowledge lives in:", base_out)
    print("Tip: Index both the .txt files and any .mermaid.mmd files for diagram-aware answers.")

if __name__ == "__main__":
    main()
