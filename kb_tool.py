#!/usr/bin/env python
# -*- coding: utf-8 -*-

r"""
kb_tool.py — One-file CLI for building a GPT-friendly knowledge base (NO Mermaid, ONE .txt per file).

What it does
- PDFs: extract text; OCR page scans; OCR embedded images; confidence-aware cleaning; add PromptSuggestion.
- PPTX: keep slide text; extract tables (TSV-like); OCR slide pictures (multi-pass); preserve tiny engineering labels.
- DOCX/XLS(X)/PNG/JPG: extract text/tables or OCR images similarly.
- Bundle: combine all .txt into one bundle per top-level folder.
- Split-text: (optional helper) split a large .txt by size.
- Doctor: quick environment checks (Tesseract).

Examples (CMD):
  python kb_tool.py convert --input C:\kb\raw --output C:\kb\output ^
    --tesseract "C:\Users\YOU\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"

  python kb_tool.py bundle --output C:\kb\output
"""

import argparse, json, math, re, sys, shutil, io, unicodedata
from pathlib import Path
from typing import List

# ---------- graceful imports ----------
try: import fitz  # PyMuPDF
except Exception: fitz = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try: import docx
except Exception: docx = None
try: import pandas as pd
except Exception: pd = None
try:
    from PIL import Image, ImageFilter, ImageOps
except Exception:
    Image = None
    ImageFilter = None
    ImageOps = None
try: import pytesseract
except Exception: pytesseract = None
try: import tiktoken
except Exception: tiktoken = None
try:
    from tqdm import tqdm
except Exception:
    def tqdm(x, **k): return x

# ---------- ignore junk/lock files ----------
IGNORE_PREFIXES = ("~$", "._")
IGNORE_BASENAMES = {"Thumbs.db", "desktop.ini"}
def is_ignored(p: Path) -> bool:
    name = p.name
    return name.startswith(IGNORE_PREFIXES) or name in IGNORE_BASENAMES

# ---------- utils ----------
def clean_ws(s: str) -> str:
    s = re.sub(r"\r\n?", "\n", s)
    s = re.sub(r"[ \t]+\n", "\n", s)
    return s.strip()

def set_tesseract(path_str: str|None):
    if pytesseract and path_str:
        p = Path(path_str)
        if p.exists():
            pytesseract.pytesseract.tesseract_cmd = str(p)

def tokenize_len(text: str) -> int:
    if tiktoken:
        enc = tiktoken.get_encoding("cl100k_base")
        return len(enc.encode(text))
    return max(1, math.ceil(len(text)/4))  # crude fallback (≈4 chars/token)

def chunk_text(text: str, max_tokens: int, overlap_tokens: int) -> List[str]:
    """Return ONE chunk unless max_tokens>0 and text exceeds that limit."""
    if max_tokens is None or max_tokens <= 0:
        return [text]  # <-- no chunking by default
    if tokenize_len(text) <= max_tokens:
        return [text]
    if tiktoken:
        enc = tiktoken.get_encoding("cl100k_base")
        toks = enc.encode(text)
        chunks, i, step = [], 0, max(1, max_tokens - overlap_tokens)
        while i < len(toks):
            chunks.append(enc.decode(toks[i:i+max_tokens]))
            i += step
        return chunks
    # char-approx fallback
    avg = 4
    max_chars = max_tokens * avg
    overlap_chars = overlap_tokens * avg
    chunks, i = [], 0
    while i < len(text):
        chunks.append(text[i:i+max_chars])
        i += max_chars - overlap_chars
    return chunks

def safe_write(path: Path, text: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8", errors="ignore")

# ---------- image description heuristics (improved) ----------
CHART_WORDS = {"chart","graph","plot","axis","axes","trend","r²","r^2","slope","time","day","month","year","kpi","rate"}
FLOW_WORDS  = {"process","flow","pid","p&id","pump","tank","filter","valve","stage","pass","membrane","ro","diagram","intake","discharge","treatment"}

UNIT_PATTERNS = [
    r"\bmg/?l\b", r"\bµs/?cm\b", r"\bus/?cm\b", r"\bppm\b", r"\bppb\b",
    r"\bmeq/?l\b", r"\bm(eq|e)/?l\b", r"\bgpm\b", r"\bm3/?h\b", r"\bm³/?h\b",
    r"\bpsi\b", r"\bbar\b", r"\bkpa\b", r"\bph\b"
]
PARAM_TOKENS = [
    r"\btds\b", r"\btss\b", r"\btoc\b", r"\balkalinity\b",
    r"\b(total|dissolved|suspended)\b", r"\bcalcium\b", r"\bmagnesium\b",
    r"\bsodium\b", r"\bpotassium\b", r"\bchloride\b", r"\bnitrate\b",
    r"\bphosphate|\bphosphorus\b", r"\bfluoride\b", r"\bstrontium\b",
    r"\bsilica|\bsio2\b"
]
UNIT_PATTERNS = [re.compile(p, re.I) for p in UNIT_PATTERNS]
PARAM_PATTERNS = [re.compile(p, re.I) for p in PARAM_TOKENS]

def _count_hits(text: str, pats):
    return sum(1 for p in pats if p.search(text))

def guess_image_type(ocr_text: str) -> str:
    """Classify OCR block: PARAM_TABLE vs CHART_OR_GRAPH vs FLOW_DIAGRAM vs GENERIC_FIGURE."""
    t = (ocr_text or "").lower()
    if not t.strip():
        return "GENERIC_FIGURE"

    pipe_count = t.count("|")  # tables often have many vertical bars
    unit_hits  = _count_hits(t, UNIT_PATTERNS)
    param_hits = _count_hits(t, PARAM_PATTERNS)

    if pipe_count >= 3 or (unit_hits + param_hits) >= 4:
        return "PARAM_TABLE"
    if "flow" in t and ("gpm" in t or "m3/h" in t or "m³/h" in t) and param_hits >= 2:
        return "PARAM_TABLE"

    c = sum(w in t for w in CHART_WORDS)
    if c >= 2:
        return "CHART_OR_GRAPH"

    f = sum(w in t for w in FLOW_WORDS)
    if f >= 2:
        return "FLOW_DIAGRAM"

    return "GENERIC_FIGURE"

def prompt_suggestion(img_type: str, ocr_text: str) -> str:
    t = (ocr_text or "")[:400]
    if img_type == "PARAM_TABLE":
        return (
            "Reconstruct a clean table with columns: Parameter | Units | Value. "
            "Parse units carefully (mg/L, µS/cm, gpm). Then, if helpful, create a simple bar chart "
            "for a subset of key parameters (e.g., TDS, TSS, TOC, major ions). "
            f"Use these lines: {t} ..."
        )
    if img_type == "CHART_OR_GRAPH":
        return (
            "Recreate a clean chart using nearby tabular values and the labels below. "
            "Choose an appropriate chart type and label axes clearly. "
            f"Labels: {t} ..."
        )
    if img_type == "FLOW_DIAGRAM":
        return (
            "Draw a simple left-to-right process flow (boxes + arrows). "
            "Use short step names extracted from text around this figure. "
            f"Helpful terms: {t} ..."
        )
    return (
        "Generate a clear illustration that communicates the key ideas. "
        "Use short labels and simple layout. "
        f"Helpful text: {t} ..."
    )

# ---------- OCR config & preprocessing ----------
OCR_CONFIG = "--oem 3 --psm 6 -c preserve_interword_spaces=1"

def preprocess_image_for_ocr(img):
    """Light, fast preprocessing that helps most scanned docs/charts."""
    try:
        g = img.convert("L")
        g = ImageOps.autocontrast(g)
        if ImageFilter:
            g = g.filter(ImageFilter.MedianFilter(size=3))
        return g
    except Exception:
        return img

# ---------- OCR cleaning & detail heuristics ----------
SAFE_ASCII = set("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789%/().,:; -_+<>={}[]|^=~")
WORD_RE = re.compile(r"\b[0-9A-Za-z][0-9A-Za-z._/%+-]*\b")

# Engineering-ish tokens common in PFDs/P&IDs/plots
ENG_PATTERNS = [
    r"\b[PV]-\d+[A-Z]?\b",                       # P-101A, V-23
    r"\b(?:FIC|TIC|LIC|PIC|PT|FT|TT|LT)[-\s]?\d+\b",
    r"\b(?:gpm|l\/s|m3\/h|m³\/h|nm3\/h|psi|bar|kpa|pa)\b",
    r"\bmg\/l\b|\bµs\/cm\b|\bus\/cm\b",          # mg/L, µS/cm
    r"\bro\b|\bnf\b|\buf\b|\bcip\b",             # process abbreviations
]
ENG_PATTERNS = [re.compile(p, re.I) for p in ENG_PATTERNS]

def has_engineering_details(s: str) -> bool:
    if not s: return False
    hits = 0
    for pat in ENG_PATTERNS:
        if pat.search(s):
            hits += 1
            if hits >= 2:
                return True
    return False

def clean_ocr_noise(raw: str):
    """
    Normalize & filter OCR text.
    Returns: (clean_text, is_low_quality, token_count, has_eng_detail)
    """
    s = raw or ""
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch if (ch in SAFE_ASCII or ch == "\n") else " " for ch in s)
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s).strip()

    tokens = WORD_RE.findall(s)
    token_count = len(tokens)
    alnum = sum(ch.isalnum() for ch in s)
    ratio = alnum / max(1, len(s))
    eng = has_engineering_details(s)

    # base low-quality gate
    low = (token_count < 4) or (ratio < 0.45)
    # override: keep if engineering detail present
    if eng and token_count >= 2:
        low = False

    return s, low, token_count, eng

# --- Confidence-aware, multi-pass OCR ---
def ocr_multipass(img, lang: str, config: str):
    """
    Confidence-aware OCR:
      - Try multiple scales + PSM modes
      - Keep only high-confidence tokens
      - Prefer results that contain engineering labels/units
    Returns: (best_clean, best_low, best_raw, eng_detail_flag)
    """
    if Image is None or pytesseract is None:
        return "", True, "", False

    try:
        from pytesseract import Output
    except Exception:
        Output = None

    base = preprocess_image_for_ocr(img)
    scales = [1.5, 2.0, 3.0] if min(base.size) < 1400 else [1.25, 1.5, 2.0]
    psms = [11, 6, 4]  # sparse text, default, block/columns

    CONF_MIN = 65     # min per-word confidence
    MAX_WORD_LEN = 30

    best = {"clean":"", "low":True, "raw":"", "tokens":0, "eng":False, "score":(0,0)}
    for sc in scales:
        w = max(1, int(base.width * sc)); h = max(1, int(base.height * sc))
        up = base.resize((w, h), resample=Image.LANCZOS)

        for psm in psms:
            try:
                if Output is not None:
                    cfg = f"--oem 3 --psm {psm} -c preserve_interword_spaces=1"
                    data = pytesseract.image_to_data(up, lang=lang, config=cfg, output_type=Output.DICT)
                    lines = {}
                    n = len(data.get("text", []))
                    for i in range(n):
                        word = (data["text"][i] or "").strip()
                        try:
                            conf = float(data["conf"][i])
                        except Exception:
                            conf = -1.0
                        if conf < CONF_MIN: 
                            continue
                        if not word or len(word) > MAX_WORD_LEN:
                            continue
                        if not all((ch.isalnum() or ch in SAFE_ASCII) for ch in word):
                            continue
                        key = (data.get("block_num",[0])[i], data.get("par_num",[0])[i], data.get("line_num",[0])[i])
                        lines.setdefault(key, []).append(word)
                    raw = "\n".join(" ".join(words) for _, words in sorted(lines.items()))
                else:
                    cfg = f"--oem 3 --psm {psm} -c preserve_interword_spaces=1"
                    raw = pytesseract.image_to_string(up, lang=lang, config=cfg)
            except Exception:
                continue

            clean, low, tokens, eng = clean_ocr_noise(raw)
            score = (1 if eng else 0, tokens)
            if score > best["score"]:
                best.update({"clean": clean, "low": low, "raw": raw, "tokens": tokens, "eng": eng, "score": score})

    return best["clean"], best["low"], best["raw"], best["eng"]

# ---------- converters ----------
def convert_pdf(path: Path, ocr_lang: str) -> str:
    out = [f"[SOURCE] {path}", ""]
    if not fitz:
        out.append("(PyMuPDF not installed)")
        return "\n".join(out)
    doc = fitz.open(path)
    for i, page in enumerate(doc, 1):
        out.append(f"\n=== PDF Page {i} ===")
        text = clean_ws(page.get_text("text") or "")
        if text:
            out.append(text)

        # OCR fallback for page scans/diagrams
        if (not text or len(text) < 30) and pytesseract and Image:
            try:
                pix = page.get_pixmap(dpi=300, alpha=False)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                clean, low, _raw, eng = ocr_multipass(img, ocr_lang, OCR_CONFIG)
                if clean and not low:
                    out.append("\n[OCR Page]")
                    out.append(clean)
                elif eng and clean:
                    out.append("\n[OCR Page (fine details)]")
                    out.append(clean)
                else:
                    out.append("\n[OCR Page skipped: low quality]")
            except Exception:
                pass

        # OCR embedded images (descriptions + prompt suggestions)
        try:
            if pytesseract and Image:
                for idx, (xref, *_rest) in enumerate(page.get_images(full=True), 1):
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n >= 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    clean, low, _raw, eng = ocr_multipass(img, ocr_lang, OCR_CONFIG)
                    itype  = guess_image_type(clean if (not low or eng) else "")
                    prompt = prompt_suggestion(itype, clean if (not low or eng) else "")
                    out.append(f"\n[IMAGE {idx} | {itype}]")
                    if clean and (not low or eng):
                        out.append(f"OCR: {clean}")
                    else:
                        out.append("OCR: (skipped: low quality)")
                    out.append(f"PromptSuggestion: {prompt}")
        except Exception:
            pass
    doc.close()
    return "\n".join(out)

def convert_image(path: Path, ocr_lang: str) -> str:
    if not (pytesseract and Image):
        return f"[SOURCE] {path}\n(OCR not available)"
    img = Image.open(path).convert("RGB")
    clean, low, _raw, eng = ocr_multipass(img, ocr_lang, OCR_CONFIG)
    itype = guess_image_type(clean if (not low or eng) else "")
    prompt = prompt_suggestion(itype, clean if (not low or eng) else "")
    if clean and (not low or eng):
        ocr_line = f"OCR: {clean}"
    else:
        ocr_line = "OCR: (skipped: low quality)"
    return f"[SOURCE] {path}\n[IMAGE | {itype}]\n{ocr_line}\nPromptSuggestion: {prompt}"

def convert_docx(path: Path) -> str:
    if not docx:
        return f"[SOURCE] {path}\n(python-docx not installed)"
    d = docx.Document(str(path))
    parts = [f"[SOURCE] {path}", ""]
    for p in d.paragraphs:
        t = p.text.strip()
        if t: parts.append(t)
    # tables
    for ti, tbl in enumerate(d.tables, 1):
        parts.append(f"\n[Table {ti}]")
        for row in tbl.rows:
            cells = [(c.text or "").replace("\n"," ").strip() for c in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)

def convert_excel(path: Path) -> str:
    if not pd:
        return f"[SOURCE] {path}\n(pandas not installed)"
    parts = [f"[SOURCE] {path}", ""]
    try:
        xls = pd.ExcelFile(path, engine=None)
    except Exception as e:
        return f"[SOURCE] {path}\n(Excel read error: {e})"
    for sheet in xls.sheet_names:
        try:
            df = xls.parse(sheet)
            parts.append(f"\n=== Sheet: {sheet} ===")
            parts.append(df.to_csv(sep="\t", index=False))
        except Exception as e:
            parts.append(f"\n=== Sheet: {sheet} ===\n(Read error: {e})")
    return "\n".join(parts)

def _pptx_extract_tables(slide):
    tables = []
    for shp in getattr(slide, "shapes", []):
        if getattr(shp, "has_table", False) and getattr(shp, "table", None):
            lines = []
            for row in shp.table.rows:
                cells = []
                for cell in row.cells:
                    t = ""
                    if getattr(cell, "text_frame", None):
                        t = "\n".join(p.text for p in cell.text_frame.paragraphs)
                    elif hasattr(cell, "text"):
                        t = cell.text
                    cells.append((t or "").replace("\n"," ").strip())
                lines.append("\t".join(cells))
            tables.append("\n".join(lines))
    return tables

def _pptx_collect_shapes_with_text(slide):
    """Collect text boxes likely to be content; skip very short/noisy text."""
    items = []
    for shp in getattr(slide, "shapes", []):
        try:
            if hasattr(shp, "text_frame") and shp.text_frame:
                txt = "\n".join(p.text for p in shp.text_frame.paragraphs).strip()
                if not txt or len(txt) < 3:
                    continue
                items.append(txt)
        except Exception:
            continue
    return items

def convert_pptx(path: Path, ocr_lang: str) -> str:
    if not Presentation:
        return f"[SOURCE] {path}\n(python-pptx not installed)"
    prs = Presentation(str(path))
    parts = [f"[SOURCE] {path}", ""]
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n=== Slide {i} ===")
        # text
        for txt in _pptx_collect_shapes_with_text(slide):
            parts.append(txt.strip())
        # tables
        for ti, tsv in enumerate(_pptx_extract_tables(slide), 1):
            parts.append(f"\n[Table {ti}]")
            parts.append(tsv)
        # pictures → multipass OCR + prompt suggestion
        if pytesseract and Image:
            for shp in slide.shapes:
                try:
                    if getattr(shp, "image", None):  # picture-like
                        img_bytes = shp.image.blob
                        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        clean, low, _raw, eng = ocr_multipass(img, ocr_lang, OCR_CONFIG)
                        itype  = guess_image_type(clean if (not low or eng) else "")
                        prompt = prompt_suggestion(itype, clean if (not low or eng) else "")
                        parts.append(f"\n[Picture on Slide {i} | {itype}]")
                        if clean and (not low or eng):
                            parts.append(f"OCR: {clean}")
                        else:
                            parts.append("OCR: (skipped: low quality)")
                        parts.append(f"PromptSuggestion: {prompt}")
                except Exception:
                    pass
    return "\n".join(parts)

# ---------- drivers ----------
SUPPORTED = {".pdf",".png",".jpg",".jpeg",".pptx",".docx",".xlsx",".xls"}

def do_convert(args):
    base_in  = Path(args.input).resolve()
    base_out = Path(args.output).resolve()
    base_out.mkdir(parents=True, exist_ok=True)
    set_tesseract(args.tesseract)

    files = [p for p in base_in.rglob("*")
             if p.is_file() and p.suffix.lower() in SUPPORTED and not is_ignored(p)]
    if not files:
        print(f"No supported files found in {base_in}")
        return

    for p in tqdm(files, desc="Converting"):
        rel = p.relative_to(base_in)
        out_dir = base_out / rel.parent
        stem = p.stem
        txt_out = out_dir / f"{stem}.txt"
        ext = p.suffix.lower()

        try:
            if ext == ".pdf":
                text = convert_pdf(p, args.ocr_lang)
            elif ext in {".png",".jpg",".jpeg"}:
                text = convert_image(p, args.ocr_lang)
            elif ext == ".docx":
                text = convert_docx(p)
            elif ext in {".xlsx",".xls"}:
                text = convert_excel(p)
            elif ext == ".pptx":
                text = convert_pptx(p, args.ocr_lang)
            else:
                text = f"[SOURCE] {p}\n(Unsupported legacy format; save as modern type and re-run.)"
        except Exception as e:
            text = f"[SOURCE] {p}\n(Extraction error: {e})"

        text = clean_ws(text)
        # NO CHUNKING by default (max_tokens=0)
        chunks = chunk_text(text, args.max_tokens, args.overlap_tokens)
        if len(chunks) == 1:
            safe_write(txt_out, chunks[0])
        else:
            for i, ch in enumerate(chunks, 1):
                safe_write(out_dir / f"{stem}.part{i:03d}.txt", ch)

    print("\nDone. Converted knowledge lives in:", base_out)
    print("Tip: index the .txt files. Use PromptSuggestion lines to reconstruct diagrams/tables/charts on demand.")

def do_bundle(args):
    base = Path(args.output).resolve()
    txts = [p for p in base.rglob("*.txt") if not p.name.startswith("bundle__")]
    if not txts:
        print(f"No .txt files under {base}")
        return
    # one bundle per top-level folder
    def key(p: Path) -> str:
        rel = p.relative_to(base)
        return rel.parts[0] if len(rel.parts)>1 else "root"

    groups = {}
    for p in txts:
        groups.setdefault(key(p), []).append(p)

    for grp, files in sorted(groups.items()):
        out = base / f"bundle__{grp}.txt"
        out.parent.mkdir(parents=True, exist_ok=True)
        count = 0
        with out.open("w", encoding="utf-8", errors="ignore") as w:
            for f in sorted(files):
                w.write(f"\n\n===== FILE: {f.relative_to(base)} =====\n\n")
                try:
                    w.write(f.read_text(encoding="utf-8", errors="ignore"))
                except Exception as e:
                    w.write(f"[READ ERROR] {e}")
                count += 1
        print(f"Wrote {out} ({count} files)")

def do_split_text(args):
    src = Path(args.file).resolve()
    if not src.exists():
        print("Not found:", src); return
    data = src.read_bytes()
    chunk = args.max_mb * 1024 * 1024
    n = max(1, math.ceil(len(data)/chunk))
    for i in range(n):
        part = data[i*chunk:(i+1)*chunk]
        out = src.parent / f"{src.stem}.part{i+1:02d}{src.suffix}"
        out.write_bytes(part)
        print(f"Wrote {out} ({len(part)/1024/1024:.1f} MB)")

def do_doctor(args):
    print("Python:", sys.version)
    if pytesseract:
        if args.tesseract:
            set_tesseract(args.tesseract)
        path = getattr(pytesseract.pytesseract, "tesseract_cmd", None)
        print("Tesseract path:", path or "(not set)")
        if path and Path(path).exists():
            try:
                ver = pytesseract.get_tesseract_version()
                print("Tesseract version:", ver)
            except Exception as e:
                print("Tesseract check failed:", e)
        else:
            print("NOTE: set --tesseract to your tesseract.exe")
    else:
        print("pytesseract not installed")

# ---------- CLI ----------
def build_parser():
    p = argparse.ArgumentParser(description="KB CLI: convert/bundle/split for GPT knowledge.")
    sub = p.add_subparsers(dest="cmd", required=True)

    pc = sub.add_parser("convert", help="Convert a folder of docs into GPT-friendly text")
    pc.add_argument("--input",  required=True, help="Input folder (e.g., C:\\kb\\raw)")
    pc.add_argument("--output", required=True, help="Output folder (e.g., C:\\kb\\output)")
    pc.add_argument("--tesseract", default=None, help="Path to tesseract.exe")
    pc.add_argument("--ocr-lang", default="eng", help="OCR language (e.g., eng or 'eng+ara')")
    pc.add_argument("--max-tokens", type=int, default=0, help="0 = no chunking (one .txt per file)")
    pc.add_argument("--overlap-tokens", type=int, default=0)
    pc.set_defaults(func=do_convert)

    pb = sub.add_parser("bundle", help="Bundle many .txt into bundle__*.txt")
    pb.add_argument("--output", required=True, help="Output folder (same as convert output)")
    pb.set_defaults(func=do_bundle)

    ps = sub.add_parser("split-text", help="Split a large text file into parts by size")
    ps.add_argument("--file", required=True, help="Path to the big .txt file")
    ps.add_argument("--max-mb", type=int, default=100, help="Max part size in MB")
    ps.set_defaults(func=do_split_text)

    pd = sub.add_parser("doctor", help="Environment checks (Tesseract etc.)")
    pd.add_argument("--tesseract", default=None)
    pd.set_defaults(func=do_doctor)

    return p

def main():
    args = build_parser().parse_args()
    args.func(args)

if __name__ == "__main__":
    main()
